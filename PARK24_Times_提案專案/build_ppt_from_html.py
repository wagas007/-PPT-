from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from playwright.sync_api import sync_playwright


ROOT = Path(__file__).resolve().parent
HTML = ROOT / "PARK24_PPT_PREVIEW_EXTERNAL.html"
OUT_DIR = ROOT / "ppt_export_slides"
PPTX = ROOT / "PARK24_Times_External_Final.pptx"
EDGE = Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe")


def render_slides() -> list[Path]:
    OUT_DIR.mkdir(exist_ok=True)
    for old in OUT_DIR.glob("slide_*.png"):
        old.unlink()

    screenshots: list[Path] = []
    with sync_playwright() as p:
        browser = p.chromium.launch(
            executable_path=str(EDGE),
            args=["--allow-file-access-from-files"],
        )
        page = browser.new_page(
            viewport={"width": 1332, "height": 900},
            device_scale_factor=2,
        )
        page.goto(HTML.as_uri(), wait_until="networkidle")
        slides = page.locator(".slide")
        count = slides.count()
        for index in range(count):
            target = OUT_DIR / f"slide_{index + 1:02d}.png"
            slides.nth(index).screenshot(path=str(target))
            screenshots.append(target)
        browser.close()
    return screenshots


def build_pptx(screenshots: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    for image in screenshots:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    if len(prs.slides) > len(screenshots):
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    prs.save(PPTX)


def main() -> None:
    screenshots = render_slides()
    build_pptx(screenshots)
    print(f"slides={len(screenshots)}")
    print(PPTX)


if __name__ == "__main__":
    main()
