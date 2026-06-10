from pathlib import Path

from PIL import Image
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches


PROJECT = Path(__file__).resolve().parents[1]
HTML = PROJECT / "04_簡版PPT" / "VIVI_商開會議_EXTERNAL.html"
OUT_DIR = PROJECT / "04_簡版PPT" / "VIVI_商開會議_PNG"
PPTX = PROJECT / "04_簡版PPT" / "VIVI_商開會議_EXTERNAL.pptx"
EDGE = Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe")


def render_slides() -> list[Path]:
    OUT_DIR.mkdir(exist_ok=True)
    for old in OUT_DIR.glob("slide_*.png"):
        old.unlink()

    screenshots = []
    with sync_playwright() as p:
        launch_args = {"args": ["--allow-file-access-from-files"]}
        if EDGE.exists():
            launch_args["executable_path"] = str(EDGE)
        browser = p.chromium.launch(**launch_args)
        page = browser.new_page(
            viewport={"width": 1280, "height": 720},
            device_scale_factor=2,
        )
        page.goto(HTML.resolve().as_uri(), wait_until="networkidle")
        slides = page.locator(".slide")
        for index in range(slides.count()):
            target = OUT_DIR / f"slide_{index + 1:02d}.png"
            slides.nth(index).screenshot(path=str(target))
            screenshots.append(target)
        browser.close()
    return screenshots


def build_pptx(screenshots: list[Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for image in screenshots:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    if len(prs.slides) > len(screenshots):
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    prs.save(PPTX)


def main() -> None:
    screenshots = render_slides()
    build_pptx(screenshots)
    dimensions = sorted({Image.open(path).size for path in screenshots})
    print(f"slides={len(screenshots)}")
    print(f"png_dimensions={dimensions}")
    print(f"png_dir={OUT_DIR}")
    print(f"pptx={PPTX}")


if __name__ == "__main__":
    main()
